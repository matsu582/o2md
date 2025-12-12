#!/usr/bin/env python3
"""
PowerPoint to Markdown Converter
PowerPointファイルをスライドごとに詳細なMarkdown形式に変換するツール

特徴:
- スライドごとに見出しを設定
- テキストは文章に、箇条書きはリストに変換
- 表がある場合はMarkdownテーブルに変換
- 図形がある場合は、そのスライドにある図形群を一つの画像として出力
- 表と図形またはテキストと図形が複合している場合は、スライドごと画像化
- グループシェイプ内のシェイプを再帰的に処理
- スライドノートを出力
- チャートデータをテーブルとして出力
- 文字の書式設定（太字、斜体、下線、取り消し線、上付き/下付き）に対応
"""

import os
import sys
import re
import tempfile
import subprocess
import shutil
import zipfile
import urllib.parse
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import List, Dict, Tuple, Optional, Any
from collections import defaultdict
import io

from utils import get_libreoffice_path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError as e:
    raise ImportError(
        "python-pptxライブラリが必要です: pip install python-pptx または uv sync を実行してください"
    ) from e

try:
    from PIL import Image
except ImportError as e:
    raise ImportError(
        "Pillowライブラリが必要です: pip install pillow または uv sync を実行してください"
    ) from e

try:
    import fitz
except ImportError as e:
    raise ImportError(
        "PyMuPDFライブラリが必要です: pip install PyMuPDF または uv sync を実行してください"
    ) from e

# 設定
LIBREOFFICE_PATH = get_libreoffice_path()

# DPI設定
DEFAULT_DPI = 300
IMAGE_QUALITY = 95



# グローバルverboseフラグ
_VERBOSE = False

def set_verbose(verbose: bool):
    """verboseモードを設定"""
    global _VERBOSE
    _VERBOSE = verbose

def is_verbose() -> bool:
    """verboseモードかどうかを返す"""
    return _VERBOSE

def debug_print(*args, **kwargs):
    """verboseモード時のみ出力するデバッグ用print"""
    if _VERBOSE:
        print(*args, **kwargs)

class PowerPointToMarkdownConverter:
    def __init__(self, pptx_file_path: str, output_dir=None, output_format='png'):
        """コンバーター初期化
        
        Args:
            pptx_file_path: 変換するPowerPointファイルのパス（.pptまたは.pptx）
            output_dir: 出力ディレクトリ（省略時は./output）
            output_format: 出力画像形式 ('png' または 'svg')
        """
        self.original_file = pptx_file_path
        self.base_name = Path(pptx_file_path).stem
        self._temp_pptx_file = None
        
        # pptファイルの場合はpptxに変換
        if pptx_file_path.lower().endswith('.ppt'):
            print(f"[INFO] .pptファイルを検出。.pptxに変換します...")
            self.pptx_file = self._convert_ppt_to_pptx(pptx_file_path)
            if not self.pptx_file:
                raise RuntimeError("pptからpptxへの変換に失敗しました")
            self._temp_pptx_file = self.pptx_file  # 後でクリーンアップするためにフラグを立てる
        else:
            self.pptx_file = pptx_file_path
        
        self.prs = Presentation(self.pptx_file)
        
        # 出力ディレクトリの設定
        if output_dir:
            self.output_dir = output_dir
        else:
            self.output_dir = os.path.join(os.getcwd(), "output")
        
        self.images_dir = os.path.join(self.output_dir, "images")
        
        # ディレクトリ作成
        os.makedirs(self.output_dir, exist_ok=True)
        os.makedirs(self.images_dir, exist_ok=True)
        
        self.markdown_lines = []
        self.image_counter = 0
        self.slide_counter = 0
        self.output_format = output_format.lower() if output_format else 'png'
        
        # 出力形式の検証
        if self.output_format not in ('png', 'svg'):
            print(f"[WARNING] 不明な出力形式 '{output_format}'。'png'を使用します。")
            self.output_format = 'png'
        
        print(f"[INFO] 出力画像形式: {self.output_format.upper()}")
    
    def convert(self) -> str:
        """メイン変換処理"""
        print(f"[INFO] PowerPoint文書変換開始: {self.pptx_file}")
        
        # ドキュメント見出し
        self.markdown_lines.append(f"# {self.base_name}")
        self.markdown_lines.append("")
        
        # 各スライドを変換
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            self.slide_counter = slide_idx
            print(f"[INFO] スライド {slide_idx}/{len(self.prs.slides)} を処理中...")
            
            try:
                self._convert_slide(slide, slide_idx)
            except Exception as e:
                print(f"[WARNING] スライド {slide_idx} の処理中にエラー: {e}")
                import traceback
                traceback.print_exc()
                continue
        
        # Markdownファイルを書き出し
        markdown_content = "\n".join(self.markdown_lines)
        output_file = os.path.join(self.output_dir, f"{self.base_name}.md")
        
        with open(output_file, "w", encoding="utf-8") as f:
            f.write(markdown_content)
        
        print(f"[SUCCESS] 変換完了: {output_file}")
        return output_file
    
    def _convert_slide(self, slide, slide_idx: int):
        """スライドを変換
        
        Args:
            slide: スライドオブジェクト
            slide_idx: スライド番号
        """
        # スライドタイトルを取得
        title = self._get_slide_title(slide)
        if title:
            self.markdown_lines.append(f"## {title}")
        else:
            self.markdown_lines.append(f"## スライド {slide_idx}")
        self.markdown_lines.append("")
        
        # スライドの内容を分析
        slide_info = self._analyze_slide(slide)
        
        # 複合スライドの判定
        has_text = slide_info['has_text']
        has_table = slide_info['has_table']
        has_shapes = slide_info['has_shapes']
        has_charts = slide_info.get('has_charts', False)
        
        # 複合スライド（テキスト/表 + 図形）またはチャートがある場合の判定
        is_complex = (has_text or has_table) and has_shapes
        needs_image = is_complex or has_charts
        
        if is_complex:
            print(f"[INFO] スライド {slide_idx}: 複合スライド検出 - テキスト/表を展開後、スライド全体を画像化")
        elif has_charts:
            print(f"[INFO] スライド {slide_idx}: チャート検出 - スライド全体を画像化（チャートデータは画像内に含まれます）")
        
        # コンテンツアイテムを順序通りに出力
        if slide_info['content_items']:
            prev_type = None
            for i, item in enumerate(slide_info['content_items']):
                # 現在のアイテムのタイプを判定
                if item.startswith('- ') or item.startswith('  - '):
                    curr_type = 'bullet'
                elif item.startswith('1. ') or item.startswith('  1. '):
                    curr_type = 'numbered'
                else:
                    curr_type = 'text'
                
                # タイプが変わった場合、またはリストの後に通常テキストが来る場合は空行を挿入
                if prev_type is not None and prev_type != curr_type:
                    self.markdown_lines.append("")
                
                self.markdown_lines.append(item)
                prev_type = curr_type
            
            # 最後に空行を追加
            self.markdown_lines.append("")
        
        # 表を出力
        if slide_info['tables']:
            for table_md in slide_info['tables']:
                self.markdown_lines.append(table_md)
                self.markdown_lines.append("")
        
        # チャートを出力（チャートデータをテーブルとして抽出）
        if slide_info['charts']:
            for chart_md in slide_info['charts']:
                self.markdown_lines.append(chart_md)
                self.markdown_lines.append("")
        
        # 図形またはチャートの処理
        if has_shapes or has_charts:
            if needs_image:
                # 複合スライドまたはチャートあり：スライド全体を画像化してテキストの最後に挿入
                print(f"[INFO] スライド {slide_idx}: スライド全体を画像化")
                self._render_slide_as_image(slide, slide_idx)
            elif has_shapes:
                # 図形のみ：図形群を画像化
                print(f"[INFO] スライド {slide_idx}: 図形のみ - 図形群を画像化")
                self._render_shapes_as_image(slide, slide_idx)
        
        # スライドノートを出力
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame is not None:
                notes_text = notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    self.markdown_lines.append("")
                    self.markdown_lines.append("### ノート:")
                    self.markdown_lines.append(notes_text)
                    self.markdown_lines.append("")
        
        # スライド間の区切り
        self.markdown_lines.append("---")
        self.markdown_lines.append("")
    
    def _get_slide_title(self, slide) -> Optional[str]:
        """スライドのタイトルを取得
        
        Args:
            slide: スライドオブジェクト
            
        Returns:
            str: タイトル文字列、存在しない場合はNone
        """
        # タイトルプレースホルダーを探す（タイプが1のプレースホルダー）
        if hasattr(slide, 'shapes'):
            for shape in slide.shapes:
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    # プレースホルダータイプ1はタイトル
                    if hasattr(shape, 'placeholder_format'):
                        if shape.placeholder_format.type == 1:  # PP_PLACEHOLDER.TITLE
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                # 複数段落がある場合は<br>で連結
                                paragraphs = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
                                if paragraphs:
                                    if len(paragraphs) > 1:
                                        return '<br>'.join(paragraphs)
                                    else:
                                        # 単一段落でも改行文字が含まれている場合は<br>に変換
                                        title = paragraphs[0]
                                        # 各種改行文字を<br>に置換（\n, \r, \x0b=垂直タブ, \x0c=改ページ）
                                        if any(c in title for c in ['\n', '\r', '\x0b', '\x0c']):
                                            title = title.replace('\r\n', '<br>').replace('\n', '<br>').replace('\r', '<br>').replace('\x0b', '<br>').replace('\x0c', '<br>')
                                        return title
        
        # タイトルプレースホルダーがない場合、最初のテキストフレームをタイトルとして使用
        if hasattr(slide, 'shapes'):
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    # 複数段落がある場合は<br>で連結
                    paragraphs = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
                    if paragraphs:
                        if len(paragraphs) > 1:
                            return '<br>'.join(paragraphs)
                        else:
                            # 単一段落でも改行文字が含まれている場合は<br>に変換
                            title = paragraphs[0]
                            # 各種改行文字を<br>に置換（\n, \r, \x0b=垂直タブ, \x0c=改ページ）
                            if any(c in title for c in ['\n', '\r', '\x0b', '\x0c']):
                                title = title.replace('\r\n', '<br>').replace('\n', '<br>').replace('\r', '<br>').replace('\x0b', '<br>').replace('\x0c', '<br>')
                            return title
        
        return None
    
    def _analyze_slide(self, slide) -> Dict[str, Any]:
        """スライドの内容を分析
        
        Returns:
            dict: スライド情報
                - has_text: テキストの有無
                - has_table: 表の有無
                - has_shapes: 図形の有無
                - content_items: 順序を保持したコンテンツアイテムのリスト（各アイテムはタイプと内容を持つ）
                - tables: 表のMarkdownリスト
                - charts: チャートのMarkdownリスト
        """
        info = {
            'has_text': False,
            'has_table': False,
            'has_shapes': False,
            'has_charts': False,
            'content_items': [],  # 順序を保持するための統一リスト
            'tables': [],
            'charts': []
        }
        
        # テキストとして処理された図形を追跡
        processed_text_shapes = set()
        
        # シェイプを座標順にソート（視覚的な読み順に近づける）
        sorted_shapes = sorted(
            slide.shapes,
            key=lambda x: (
                float('inf') if x.top is None else x.top,
                float('inf') if x.left is None else x.left,
            ),
        )
        
        for shape in sorted_shapes:
            # プレースホルダーのチェック（タイトル、フッター、スライド番号などを除外）
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                if hasattr(shape, 'placeholder_format'):
                    placeholder_type = shape.placeholder_format.type
                    # 除外するプレースホルダータイプ: 1=タイトル, 12=フッター, 13=スライド番号, 14=日付
                    if placeholder_type in [1, 12, 13, 14]:
                        continue
            

            # テキストボックスまたはプレースホルダー
            if shape.has_text_frame:
                text_frame = shape.text_frame
                
                # リストタイプを判定
                list_type = self._get_list_type(text_frame)
                
                # テキストとして実際に処理した場合のみマーク
                text_was_processed = False
                
                if list_type == 'bullet':
                    # 箇条書きリスト（書式付き）
                    for paragraph in text_frame.paragraphs:
                        formatted_text = self._get_formatted_paragraph_text(paragraph)
                        if formatted_text.strip():
                            level = paragraph.level
                            indent = "  " * level
                            info['content_items'].append(f"{indent}- {formatted_text.strip()}")
                            text_was_processed = True
                    info['has_text'] = True
                elif list_type == 'numbered':
                    # 番号付きリスト（書式付き）
                    for paragraph in text_frame.paragraphs:
                        formatted_text = self._get_formatted_paragraph_text(paragraph)
                        if formatted_text.strip():
                            # 先頭の番号記号を削除
                            cleaned_text = self._remove_number_prefix(formatted_text.strip())
                            level = paragraph.level
                            indent = "  " * level
                            info['content_items'].append(f"{indent}1. {cleaned_text}")
                            text_was_processed = True
                    info['has_text'] = True
                else:
                    # 通常のテキスト: 段落ごとに個別に判定（書式付き）
                    paragraphs = []
                    for paragraph in text_frame.paragraphs:
                        plain_text = paragraph.text.strip()
                        if not plain_text:
                            continue
                        
                        # 書式付きテキストを取得
                        formatted_text = self._get_formatted_paragraph_text(paragraph)
                        
                        # 元のテキスト（インデント情報を保持）
                        original_text = paragraph.text
                        
                        # 各段落のテキストパターンを個別に判定（プレーンテキストでパターン判定）
                        if self._is_numbered_text([plain_text]):
                            # 番号付きリスト項目
                            cleaned_text = self._remove_number_prefix(formatted_text.strip())
                            info['content_items'].append(f"1. {cleaned_text}")
                            info['has_text'] = True
                            text_was_processed = True
                        elif plain_text.startswith('・') or plain_text.startswith('•'):
                            # 箇条書き項目（・や•で始まる）
                            bullet_text = formatted_text.strip().lstrip('・•').strip()
                            info['content_items'].append(f"- {bullet_text}")
                            info['has_text'] = True
                            text_was_processed = True
                        elif plain_text.startswith('-') or plain_text.startswith('−'):
                            # ハイフンやマイナスで始まる箇条書き項目、インデント検出
                            indent_match = len(original_text) - len(original_text.lstrip('　 '))
                            indent_level = indent_match // 2  # 2文字で1レベルとする
                            indent = "  " * indent_level
                            bullet_text = formatted_text.strip().lstrip('-−').strip()
                            info['content_items'].append(f"{indent}- {bullet_text}")
                            info['has_text'] = True
                            text_was_processed = True
                        else:
                            # 通常のテキスト: 一時的にリストに追加
                            paragraphs.append(formatted_text.strip() if formatted_text.strip() else plain_text)
                    
                    # 通常のテキストをまとめて処理（改行を<br>で表現）
                    if paragraphs:
                        # 複数の段落がある場合は<br>で連結
                        if len(paragraphs) > 1:
                            combined_text = '<br>'.join(paragraphs)
                            info['content_items'].append(combined_text)
                        else:
                            info['content_items'].append(paragraphs[0])
                        info['has_text'] = True
                        text_was_processed = True
                
                # テキストとして処理された図形をマーク
                # PLACEHOLDERまたはTEXT_BOXのみテキストとしてマーク
                # AUTO_SHAPEはテキストがあっても図形として扱う（ユーザー要件）
                should_mark_as_text = False
                
                if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    # プレースホルダは常にテキストとして扱う
                    should_mark_as_text = True
                elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    # テキストボックスはテキストとして扱う
                    should_mark_as_text = True
                # AUTO_SHAPE、LINE等はテキストがあっても図形として扱う
                
                if should_mark_as_text:
                    processed_text_shapes.add(id(shape))
            
            # 表
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_md = self._convert_table_shape(shape)
                if table_md:
                    info['tables'].append(table_md)
                info['has_table'] = True
            
            # チャートの処理
            if shape.has_chart:
                chart_md = self._convert_chart_to_markdown(shape.chart)
                if chart_md:
                    info['charts'].append(chart_md)
                info['has_charts'] = True
                # チャートはデータを抽出したので図形としては扱わない
                continue
            
            # グループシェイプの再帰処理
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_info = self._process_group_shape(shape, processed_text_shapes)
                info['has_text'] = info['has_text'] or group_info['has_text']
                info['has_table'] = info['has_table'] or group_info['has_table']
                info['has_shapes'] = info['has_shapes'] or group_info['has_shapes']
                info['has_charts'] = info['has_charts'] or group_info.get('has_charts', False)
                info['content_items'].extend(group_info['content_items'])
                info['tables'].extend(group_info['tables'])
                info['charts'].extend(group_info['charts'])
                continue
            
            # 図形（AutoShape, Picture など）
            if shape.shape_type in [
                MSO_SHAPE_TYPE.AUTO_SHAPE,
                MSO_SHAPE_TYPE.PICTURE,
                MSO_SHAPE_TYPE.FREEFORM,
                MSO_SHAPE_TYPE.LINE,
            ]:
                # テキストとして既に処理された図形は除外
                if id(shape) in processed_text_shapes:
                    continue
                
                # 小さい装飾図形を除外（幅と高さが両方とも1.5cm以下）
                if hasattr(shape, 'width') and hasattr(shape, 'height'):
                    small_size_emu = 540000  # 1.5cm (EMU単位)
                    if shape.width <= small_size_emu and shape.height <= small_size_emu:
                        continue
                
                # 図形として扱うかどうかを判定
                # PICTURE, FREEFORM, AUTO_SHAPE, LINEは
                # テキストボックスやプレースホルダとは明確に異なるため、
                # 塗りつぶしや枠線の有無に関わらず図形として扱う
                # （テキストボックスはshape_type=TEXT_BOXで既に除外されている）
                info['has_shapes'] = True
        
        return info
    
    def _process_group_shape(self, group_shape, processed_text_shapes: set) -> Dict[str, Any]:
        """グループシェイプ内のシェイプを再帰的に処理
        
        Args:
            group_shape: グループシェイプオブジェクト
            processed_text_shapes: テキストとして処理済みのシェイプIDセット
            
        Returns:
            dict: グループ内のコンテンツ情報
        """
        info = {
            'has_text': False,
            'has_table': False,
            'has_shapes': False,
            'has_charts': False,
            'content_items': [],
            'tables': [],
            'charts': []
        }
        
        # グループ内のシェイプを座標順にソート
        sorted_shapes = sorted(
            group_shape.shapes,
            key=lambda x: (
                float('inf') if x.top is None else x.top,
                float('inf') if x.left is None else x.left,
            ),
        )
        
        for shape in sorted_shapes:
            # テキストフレームを持つシェイプ
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    # 書式付きテキストを取得
                    formatted_text = self._get_formatted_text_from_text_frame(shape.text_frame)
                    if formatted_text:
                        info['content_items'].append(formatted_text)
                        info['has_text'] = True
                        processed_text_shapes.add(id(shape))
            
            # 表
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_md = self._convert_table_shape(shape)
                if table_md:
                    info['tables'].append(table_md)
                info['has_table'] = True
            
            # チャート
            if hasattr(shape, 'has_chart') and shape.has_chart:
                chart_md = self._convert_chart_to_markdown(shape.chart)
                if chart_md:
                    info['charts'].append(chart_md)
                info['has_charts'] = True
                continue
            
            # ネストされたグループシェイプ（再帰処理）
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                nested_info = self._process_group_shape(shape, processed_text_shapes)
                info['has_text'] = info['has_text'] or nested_info['has_text']
                info['has_table'] = info['has_table'] or nested_info['has_table']
                info['has_shapes'] = info['has_shapes'] or nested_info['has_shapes']
                info['has_charts'] = info['has_charts'] or nested_info.get('has_charts', False)
                info['content_items'].extend(nested_info['content_items'])
                info['tables'].extend(nested_info['tables'])
                info['charts'].extend(nested_info['charts'])
                continue
            
            # その他の図形
            if shape.shape_type in [
                MSO_SHAPE_TYPE.AUTO_SHAPE,
                MSO_SHAPE_TYPE.PICTURE,
                MSO_SHAPE_TYPE.FREEFORM,
                MSO_SHAPE_TYPE.LINE,
            ]:
                if id(shape) not in processed_text_shapes:
                    info['has_shapes'] = True
        
        return info
    
    def _convert_chart_to_markdown(self, chart) -> Optional[str]:
        """チャートデータをMarkdownテーブルに変換
        
        Args:
            chart: チャートオブジェクト
            
        Returns:
            str: Markdownテーブル文字列、変換できない場合はNone
        """
        try:
            md_lines = []
            
            # チャートタイトル
            md_lines.append("### チャート")
            if chart.has_title:
                try:
                    title_text = chart.chart_title.text_frame.text
                    if title_text:
                        md_lines[-1] = f"### チャート: {title_text}"
                except Exception:
                    pass
            md_lines.append("")
            
            # カテゴリと系列データを取得
            try:
                category_names = [c.label for c in chart.plots[0].categories]
                series_names = [s.name for s in chart.series]
                
                # ヘッダー行
                header = ["カテゴリ"] + series_names
                md_lines.append("| " + " | ".join(str(h) if h else "" for h in header) + " |")
                md_lines.append("|" + "|".join(["---"] * len(header)) + "|")
                
                # データ行
                for idx, category in enumerate(category_names):
                    row = [str(category) if category else ""]
                    for series in chart.series:
                        try:
                            value = series.values[idx]
                            row.append(str(value) if value is not None else "")
                        except (IndexError, TypeError):
                            row.append("")
                    md_lines.append("| " + " | ".join(row) + " |")
                
                return "\n".join(md_lines)
                
            except (ValueError, AttributeError) as e:
                # サポートされていないチャートタイプ
                if "unsupported plot type" in str(e):
                    return "### チャート\n\n[サポートされていないチャートタイプ]"
                return None
                
        except Exception as e:
            debug_print(f"[DEBUG] チャート変換エラー: {e}")
            return None
    
    def _get_formatted_text_from_text_frame(self, text_frame) -> str:
        """テキストフレームから書式付きテキストを取得
        
        Args:
            text_frame: テキストフレームオブジェクト
            
        Returns:
            str: 書式付きテキスト
        """
        result_parts = []
        
        for paragraph in text_frame.paragraphs:
            para_text = self._get_formatted_paragraph_text(paragraph)
            if para_text.strip():
                result_parts.append(para_text)
        
        return "<br>".join(result_parts) if result_parts else ""
    
    def _get_formatted_paragraph_text(self, paragraph) -> str:
        """段落から書式付きテキストを取得
        
        Args:
            paragraph: 段落オブジェクト
            
        Returns:
            str: 書式付きテキスト
        """
        para_parts = []
        for run in paragraph.runs:
            text = run.text
            if text:
                formatted_text = self._apply_run_formatting(run, text)
                para_parts.append(formatted_text)
        
        return "".join(para_parts)
    
    def _apply_run_formatting(self, run, text: str) -> str:
        """Runのフォーマット情報をMarkdown記法に変換
        
        python-pptxのFontオブジェクトで利用可能な属性:
        - bold: 太字
        - italic: 斜体
        - underline: 下線
        
        注: superscript, subscript, strikethroughはpython-pptxでは
        直接サポートされていないため、XMLから取得する必要がある
        
        Args:
            run: python-pptxのRunオブジェクト
            text: 元のテキスト
            
        Returns:
            str: フォーマット適用後のテキスト
        """
        if not text or not text.strip():
            return text
        
        # XMLから追加の書式情報を取得
        try:
            rPr = run._r.get_or_add_rPr()
            
            # 上付き文字（baseline属性で判定）
            baseline = rPr.attrib.get(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}baseline'
            )
            if baseline and int(baseline) > 0:
                text = f"<sup>{text}</sup>"
            elif baseline and int(baseline) < 0:
                # 下付き文字
                text = f"<sub>{text}</sub>"
            
            # 取り消し線
            strike_elem = rPr.find(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}strike'
            )
            if strike_elem is not None:
                strike_val = strike_elem.attrib.get('val', 'sngStrike')
                if strike_val != 'noStrike':
                    text = f"~~{text}~~"
        except Exception:
            pass
        
        # 下線（MarkdownではHTMLタグを使用）
        if run.font.underline and run.font.underline is not False:
            text = f"<u>{text}</u>"
        
        # 斜体
        if run.font.italic:
            text = f"*{text}*"
        
        # 太字
        if run.font.bold:
            text = f"**{text}**"
        
        return text
    
    def _get_list_type(self, text_frame) -> Optional[str]:
        """テキストフレームのリストタイプを判定
        
        Args:
            text_frame: テキストフレームオブジェクト
            
        Returns:
            str: 'bullet'（箇条書き）、'numbered'（番号付き）、None（通常テキスト）
        """
        for para in text_frame.paragraphs:
            if not para.text.strip():
                continue
            
            # インデントレベルがある場合はリスト
            if para.level > 0:
                # XMLから番号付きか箇条書きかを判定
                try:
                    pPr = para._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
                    if pPr is not None:
                        # buAutoNum（自動番号付け）があれば番号付きリスト
                        buAutoNum = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum')
                        if buAutoNum is not None:
                            return 'numbered'
                        
                        # buChar（箇条書き文字）があれば箇条書き
                        buChar = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
                        if buChar is not None:
                            return 'bullet'
                except Exception:
                    pass
                
                # デフォルトは箇条書き
                return 'bullet'
            
            # XMLから書式情報を取得
            try:
                pPr = para._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
                if pPr is not None:
                    # buAutoNum（自動番号付け）があれば番号付きリスト
                    buAutoNum = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum')
                    if buAutoNum is not None:
                        return 'numbered'
                    
                    # buChar（箇条書き文字）があれば箇条書き
                    buChar = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
                    if buChar is not None:
                        return 'bullet'
            except Exception:
                pass
        
        return None
    
    def _is_numbered_text(self, paragraphs: List[str]) -> bool:
        """テキストが番号付きリストかどうかを判定（テキストパターンから）
        
        Args:
            paragraphs: 段落のリスト
            
        Returns:
            bool: 番号付きリストの場合True
        """
        if not paragraphs:
            return False
        
        # 番号パターン（１．、1.、①、(1)など）
        # 注: 1-8 のようなハイフン付き番号は除外（通常のテキストとして扱う）
        import re
        patterns = [
            r'^[０-９]{1,2}[\.．、](?!\d)',  # 全角数字 + 記号（後ろに数字が続かない）
            r'^[0-9]{1,2}[\.．、](?!\d)',    # 半角数字 + 記号（後ろに数字が続かない）
            r'^[①-⑳]',                       # 丸数字
            r'^\([0-9]{1,2}\)',              # (1)形式
            r'^\（[0-9]{1,2}\）',            # （１）形式
        ]
        
        # 単一段落の場合: 番号パターンに一致すればTrue
        if len(paragraphs) == 1:
            text = paragraphs[0]
            return any(re.match(pattern, text) for pattern in patterns)
        
        # 複数段落の場合: 少なくとも2つの段落が番号パターンに一致すればTrue
        matches = 0
        for para in paragraphs[:5]:  # 最初の5段落まで確認
            for pattern in patterns:
                if re.match(pattern, para):
                    matches += 1
                    break
        
        return matches >= 2
    
    def _remove_number_prefix(self, text: str) -> str:
        """テキストの先頭から番号接頭辞を削除する
        
        Args:
            text: 処理対象のテキスト
            
        Returns:
            番号接頭辞を削除したテキスト
        """
        import re
        # 番号パターン（全角数字、半角数字、丸数字、括弧付き数字）
        patterns = [
            r'^[０-９]{1,2}[\.．、]\s*',  # 全角数字
            r'^[0-9]{1,2}[\.．、]\s*',    # 半角数字
            r'^[①-⑳]\s*',                 # 丸数字
            r'^\([0-9]{1,2}\)\s*',        # 半角括弧
            r'^\([０-９]{1,2}\)\s*',      # 全角括弧
            r'^\（[0-9]{1,2}\）\s*',      # 全角括弧
            r'^\（[０-９]{1,2}\）\s*',    # 全角括弧
        ]
        
        for pattern in patterns:
            text = re.sub(pattern, '', text)
        
        return text
    
    def _get_slide_alt_text(self, slide, slide_idx: int) -> str:
        """スライドのaltテキストを取得
        
        スライドのタイトルまたは説明を取得してaltテキストとして返す。
        タイトルがない場合は「スライド N」を返す。
        
        Args:
            slide: スライドオブジェクト
            slide_idx: スライド番号
            
        Returns:
            str: altテキスト
        """
        # スライドタイトルを取得
        title = self._get_slide_title(slide)
        if title:
            # HTMLタグを除去してaltテキストに使用
            clean_title = title.replace('<br>', ' ').strip()
            if clean_title:
                return f"スライド {slide_idx}: {clean_title}"
        
        return f"スライド {slide_idx}"
    
    def _convert_table_shape(self, shape) -> Optional[str]:
        """表シェイプをMarkdownテーブルに変換
        
        Args:
            shape: 表シェイプオブジェクト
            
        Returns:
            str: Markdownテーブル文字列
        """
        try:
            table = shape.table
            
            if not table:
                return None
            
            # 行数と列数を取得
            row_count = len(table.rows)
            col_count = len(table.columns)
            
            if row_count == 0 or col_count == 0:
                return None
            
            md_lines = []
            
            # ヘッダー行（最初の行）
            header_cells = []
            for col_idx in range(col_count):
                try:
                    cell = table.cell(0, col_idx)
                    text = cell.text.strip().replace('\n', '<br>')
                    header_cells.append(text if text else " ")
                except Exception:
                    header_cells.append(" ")
            
            md_lines.append("| " + " | ".join(header_cells) + " |")
            md_lines.append("| " + " | ".join(["---"] * len(header_cells)) + " |")
            
            # データ行
            for row_idx in range(1, row_count):
                cells = []
                for col_idx in range(col_count):
                    try:
                        cell = table.cell(row_idx, col_idx)
                        text = cell.text.strip().replace('\n', '<br>')
                        # パイプ文字をエスケープ
                        text = text.replace('|', '\\|')
                        cells.append(text if text else " ")
                    except Exception:
                        cells.append(" ")
                
                md_lines.append("| " + " | ".join(cells) + " |")
            
            return "\n".join(md_lines)
            
        except Exception as e:
            print(f"[WARNING] 表の変換エラー: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def _render_slide_as_image(self, slide, slide_idx: int):
        """スライド全体を画像として出力（元のPowerPointファイルから直接変換）
        
        Args:
            slide: スライドオブジェクト
            slide_idx: スライド番号
        """
        try:
            # 元のPowerPointファイルをPDFに変換（全スライド）
            pdf_path = self._get_or_create_pdf()
            if not pdf_path:
                return
            
            # PDFから該当ページ（スライド）を画像に変換
            self.image_counter += 1
            ext = self.output_format
            image_filename = f"{self.base_name}_slide_{slide_idx:03d}.{ext}"
            image_path = os.path.join(self.images_dir, image_filename)
            
            # PDFの該当ページ（スライドインデックスは0から始まる）を画像に変換
            success, actual_path = self._convert_pdf_page_to_image(pdf_path, slide_idx - 1, image_path)
            if success:
                # Markdownに追加（altテキストを改善）
                actual_filename = os.path.basename(actual_path)
                encoded_filename = urllib.parse.quote(actual_filename)
                alt_text = self._get_slide_alt_text(slide, slide_idx)
                self.markdown_lines.append(f"![{alt_text}](images/{encoded_filename})")
                self.markdown_lines.append("")
                print(f"[SUCCESS] スライド全体画像化: {actual_filename}")
            
        except Exception as e:
            print(f"[ERROR] スライド画像化エラー: {e}")
            import traceback
            traceback.print_exc()
    
    def _get_or_create_pdf(self) -> Optional[str]:
        """PowerPointファイル全体のPDFを取得（キャッシュ使用）
        
        Returns:
            str: PDFファイルのパス
        """
        # 既にPDFが作成されている場合はそれを返す
        if hasattr(self, '_cached_pdf_path') and os.path.exists(self._cached_pdf_path):
            return self._cached_pdf_path
        
        try:
            temp_dir = tempfile.mkdtemp()
            
            # LibreOfficeでPDFに変換
            cmd = [
                LIBREOFFICE_PATH,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', temp_dir,
                self.pptx_file
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=180)
            
            if result.returncode == 0:
                # 変換されたPDFを探す
                for file in os.listdir(temp_dir):
                    if file.endswith('.pdf'):
                        pdf_path = os.path.join(temp_dir, file)
                        # 永続的な場所にコピー
                        final_pdf_path = os.path.join(temp_dir, 'presentation.pdf')
                        shutil.copy2(pdf_path, final_pdf_path)
                        self._cached_pdf_path = final_pdf_path
                        self._temp_pdf_dir = temp_dir
                        return final_pdf_path
            
            shutil.rmtree(temp_dir)
            print(f"[ERROR] PDF変換失敗: {result.stderr}")
            return None
            
        except Exception as e:
            print(f"[ERROR] PDF変換エラー: {e}")
            return None
    
    def _convert_pdf_page_to_png(self, pdf_path: str, page_index: int, output_path: str) -> bool:
        """PDFの特定ページをPNGに変換（PyMuPDF使用）
        
        Args:
            pdf_path: PDFファイルのパス
            page_index: ページインデックス（0から始まる）
            output_path: 出力PNGファイルのパス
            
        Returns:
            bool: 変換成功時True
        """
        try:
            debug_print(f"[DEBUG] PyMuPDFでPDF→PNG変換実行 (ページ {page_index})...")
            
            doc = fitz.open(pdf_path)
            if page_index >= len(doc):
                print(f"[ERROR] ページ{page_index}が存在しません（全{len(doc)}ページ）")
                doc.close()
                return False
            
            page = doc[page_index]
            
            mat = fitz.Matrix(DEFAULT_DPI/72, DEFAULT_DPI/72)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            
            img_data = pix.tobytes("png")
            pix = None
            doc.close()
            
            img = Image.open(io.BytesIO(img_data))
            
            if img.mode == 'RGBA':
                background = Image.new('RGB', img.size, (255, 255, 255))
                background.paste(img, mask=img.split()[3] if len(img.split()) > 3 else None)
                img = background
            elif img.mode != 'RGB':
                img = img.convert('RGB')
            
            img.save(output_path, 'PNG', quality=IMAGE_QUALITY)
            
            print(f"[INFO] PNG変換完了: {output_path} (サイズ: {img.size[0]}x{img.size[1]})")
            return True
                
        except Exception as e:
            print(f"[ERROR] PNG変換エラー: {e}")
            import traceback
            traceback.print_exc()
            return False
    
    def _convert_pdf_page_to_svg(self, pdf_path: str, page_index: int, output_path: str) -> bool:
        """PDFの特定ページをSVGに変換（PyMuPDF使用）
        
        Args:
            pdf_path: PDFファイルのパス
            page_index: ページインデックス（0から始まる）
            output_path: 出力SVGファイルのパス
            
        Returns:
            bool: 変換成功時True
        """
        try:
            debug_print(f"[DEBUG] PyMuPDFでPDF→SVG変換実行 (ページ {page_index})...")
            
            doc = fitz.open(pdf_path)
            if page_index >= len(doc):
                print(f"[ERROR] ページ{page_index}が存在しません（全{len(doc)}ページ）")
                doc.close()
                return False
            
            page = doc[page_index]
            
            # SVGとして出力
            svg_content = page.get_svg_image()
            doc.close()
            
            # SVGファイルに書き込み
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(svg_content)
            
            print(f"[INFO] SVG変換完了: {output_path}")
            return True
                
        except Exception as e:
            print(f"[ERROR] SVG変換エラー: {e}")
            import traceback
            traceback.print_exc()
            return False
    
    def _convert_pdf_page_to_image(self, pdf_path: str, page_index: int, output_path: str) -> Tuple[bool, str]:
        """PDFの特定ページを画像に変換（出力形式に応じてPNGまたはSVG）
        
        Args:
            pdf_path: PDFファイルのパス
            page_index: ページインデックス（0から始まる）
            output_path: 出力ファイルのパス（PNG形式で指定）
            
        Returns:
            Tuple[bool, str]: (変換成功フラグ, 実際の出力パス)
        """
        if self.output_format == 'svg':
            # SVG出力の場合は拡張子を変更
            svg_path = output_path.replace('.png', '.svg')
            return self._convert_pdf_page_to_svg(pdf_path, page_index, svg_path), svg_path
        else:
            return self._convert_pdf_page_to_png(pdf_path, page_index, output_path), output_path
    
    def cleanup(self):
        """一時ファイルをクリーンアップ"""
        if hasattr(self, '_temp_pdf_dir') and os.path.exists(self._temp_pdf_dir):
            try:
                shutil.rmtree(self._temp_pdf_dir)
            except Exception as e:
                print(f"[WARNING] 一時ファイル削除エラー: {e}")
        
        # pptから変換された一時pptxファイルを削除
        if self._temp_pptx_file and os.path.exists(self._temp_pptx_file):
            try:
                os.remove(self._temp_pptx_file)
                print(f"[INFO] 一時pptxファイルを削除: {self._temp_pptx_file}")
            except Exception as e:
                print(f"[WARNING] 一時pptxファイル削除エラー: {e}")
    
    def _convert_ppt_to_pptx(self, ppt_file: str) -> Optional[str]:
        """pptファイルをpptxに変換
        
        Args:
            ppt_file: 変換するpptファイルのパス
            
        Returns:
            str: 変換されたpptxファイルのパス、失敗時はNone
        """
        try:
            # 一時ディレクトリを作成
            temp_dir = tempfile.mkdtemp()
            
            # LibreOfficeでpptをpptxに変換
            cmd = [
                LIBREOFFICE_PATH,
                '--headless',
                '--convert-to', 'pptx',
                '--outdir', temp_dir,
                ppt_file
            ]
            
            print(f"[INFO] LibreOfficeを使用してppt→pptx変換を実行中...")
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=180)
            
            if result.returncode == 0:
                # 変換されたpptxファイルを探す
                for file in os.listdir(temp_dir):
                    if file.endswith('.pptx'):
                        pptx_path = os.path.join(temp_dir, file)
                        print(f"[SUCCESS] ppt→pptx変換完了: {pptx_path}")
                        return pptx_path
                
                print(f"[ERROR] pptx変換失敗: 出力ファイルが見つかりません")
                shutil.rmtree(temp_dir)
                return None
            else:
                print(f"[ERROR] pptx変換失敗: {result.stderr}")
                shutil.rmtree(temp_dir)
                return None
                
        except Exception as e:
            print(f"[ERROR] pptx変換エラー: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def _render_shapes_as_image(self, slide, slide_idx: int):
        """スライド内の図形群を画像として出力
        
        Args:
            slide: スライドオブジェクト
            slide_idx: スライド番号
        """
        self._render_slide_as_image(slide, slide_idx)


def main():
    """メイン関数"""
    import argparse
    
    parser = argparse.ArgumentParser(description='PowerPoint文書をMarkdownに変換')
    parser.add_argument('pptx_file', help='変換するPowerPointファイル（.pptまたは.pptx）')
    parser.add_argument('-o', '--output-dir', type=str,
                       help='出力ディレクトリを指定（デフォルト: ./output）')
    parser.add_argument('--format', choices=['png', 'svg'], default='svg',
                       help='出力画像形式を指定（デフォルト: png）')
    parser.add_argument('-v', '--verbose', action='store_true',
                       help='デバッグ情報を出力')
    
    args = parser.parse_args()
    
    set_verbose(args.verbose)
    
    if not os.path.exists(args.pptx_file):
        print(f"エラー: ファイル '{args.pptx_file}' が見つかりません。")
        sys.exit(1)
    
    if not (args.pptx_file.lower().endswith('.pptx') or args.pptx_file.lower().endswith('.ppt')):
        print("エラー: .pptまたは.pptx形式のファイルを指定してください。")
        sys.exit(1)
    
    converter = None
    try:
        converter = PowerPointToMarkdownConverter(
            args.pptx_file,
            output_dir=args.output_dir,
            output_format=args.format
        )
        output_file = converter.convert()
        print("\n変換完了!")
        print(f"出力ファイル: {output_file}")
        print(f"画像フォルダ: {converter.images_dir}")
        
    except Exception as e:
        print(f"変換エラー: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
    finally:
        if converter:
            converter.cleanup()


if __name__ == "__main__":
    main()
